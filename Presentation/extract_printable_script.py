"""
Extract on-slide text and speaker notes from BEI_Presentation.pptx into a
printable script (slide number + presenter + content).

Run:  python Presentation/extract_printable_script.py
"""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
PPTX = ROOT / "Presentation" / "BEI_Presentation.pptx"
OUT = ROOT / "Presentation" / "BEI_Printable_Script.md"


def _collect_shape_text(shape, acc: list[str]) -> None:
    if hasattr(shape, "text"):
        t = (shape.text or "").strip()
        if t:
            acc.append(t)
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            _collect_shape_text(child, acc)


def slide_visible_text(slide) -> list[str]:
    lines: list[str] = []
    for sh in slide.shapes:
        _collect_shape_text(sh, lines)
    # De-dupe while preserving order (some decks repeat)
    seen: set[str] = set()
    out: list[str] = []
    for t in lines:
        if t not in seen:
            seen.add(t)
            out.append(t)
    return out


def _format_presenter(tag: str) -> str:
    t = tag.strip().replace("_", " ")
    return t.title() if t.isupper() or t.islower() else t


def parse_notes_header(notes: str) -> tuple[str, str]:
    """Return (presenter_label, timing_line) from first line of notes."""
    notes = (notes or "").strip()
    if not notes:
        return "—", ""
    first_line = notes.split("\n", 1)[0].strip()
    presenter = "—"
    m = re.match(r"^\[([^\]]+)\]", first_line)
    if m:
        presenter = _format_presenter(m.group(1))
    timing = ""
    m2 = re.search(r"\[(\d+:\d+)\s*-\s*(\d+:\d+)\]\s*(\d+s)?", first_line)
    if m2:
        timing = f"{m2.group(1)} – {m2.group(2)}"
        if m2.group(3):
            timing += f" ({m2.group(3)})"
    return presenter, timing


def notes_body(notes: str) -> str:
    notes = (notes or "").strip()
    if not notes:
        return ""
    parts = notes.split("\n", 1)
    return parts[1].strip() if len(parts) > 1 else ""


def main() -> None:
    prs = Presentation(str(PPTX))
    chunks: list[str] = []
    chunks.append("# Burn Equity Index — Printable Script\n")
    chunks.append(
        "_Generated from `Presentation/BEI_Presentation.pptx`. "
        "Presenter and timing are taken from the first line of speaker notes._\n"
    )

    for i, slide in enumerate(prs.slides, 1):
        visible = slide_visible_text(slide)
        ns = slide.notes_slide
        raw_notes = ""
        if ns and ns.notes_text_frame and ns.notes_text_frame.text:
            raw_notes = ns.notes_text_frame.text
        presenter, timing = parse_notes_header(raw_notes)
        body = notes_body(raw_notes)

        chunks.append(f"\n---\n\n## Slide {i}\n\n")
        chunks.append(f"**Presenter:** {presenter}\n")
        if timing:
            chunks.append(f"**Timing (from notes):** {timing}\n")
        chunks.append("\n### On-slide text\n\n")
        if visible:
            for block in visible:
                # Markdown-friendly: quote multi-line blocks
                if "\n" in block:
                    chunks.append("```\n")
                    chunks.append(block)
                    chunks.append("\n```\n\n")
                else:
                    chunks.append(f"- {block}\n")
        else:
            chunks.append("_No text shapes found._\n")

        chunks.append("\n### Speaker notes\n\n")
        if body:
            chunks.append(body)
            chunks.append("\n")
        else:
            chunks.append("_No notes._\n")

    OUT.write_text("".join(chunks), encoding="utf-8")
    print(f"Wrote {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
