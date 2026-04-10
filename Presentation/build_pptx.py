"""
Build a competition-ready PPTX presentation for the Burn Equity Index.
Colorblind-accessible, auditorium-optimized, story-driven with embedded
10:59 minute script in speaker notes.

Run:  python Presentation/build_pptx.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
DASH = ROOT / "outputs" / "dashboard"
OUT  = ROOT / "Presentation" / "BEI_Presentation.pptx"

# ═══════════════════════════════════════════════════════════════════════
#  COLORBLIND-SAFE PALETTE  (Wong / IBM adapted for dark theme)
#  Passes Deuteranopia, Protanopia, Tritanopia simulations.
# ═══════════════════════════════════════════════════════════════════════
BG         = RGBColor(0x1A, 0x23, 0x32)
CARD_BG    = RGBColor(0x24, 0x34, 0x47)
TEXT_CLR   = RGBColor(0xF0, 0xF3, 0xF7)
MUTED      = RGBColor(0x9E, 0xAD, 0xBD)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG2   = RGBColor(0x12, 0x1A, 0x27)

SKY_BLUE   = RGBColor(0x56, 0xB4, 0xE9)
AMBER      = RGBColor(0xE6, 0x9F, 0x00)
TEAL       = RGBColor(0x00, 0x9E, 0x73)
VERMILLION = RGBColor(0xD5, 0x5E, 0x00)
PINK       = RGBColor(0xCC, 0x79, 0xA7)

PILLAR_S = AMBER
PILLAR_T = VERMILLION
PILLAR_P = TEAL
PILLAR_C = SKY_BLUE

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT    = "Arial"

# Max safe content area: 0.4" margin all sides
MAX_X = 12.9   # right edge
MAX_Y = 7.1    # bottom edge


# ═══════════════════════════════════════════════════════════════════════
#  HELPERS
# ═══════════════════════════════════════════════════════════════════════

def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def tb(slide, left, top, width, height, text, size=20,
       bold=False, color=TEXT_CLR, align=PP_ALIGN.LEFT,
       font=FONT, lsp=None):
    """Add a textbox with proper line breaks and no auto-grow."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    # Disable auto-fit so shapes never grow past their bounds
    try:
        bodyPr = tf._txBody.bodyPr
        for child in list(bodyPr):
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag in ('spAutoFit', 'normAutofit', 'noAutofit'):
                bodyPr.remove(child)
        from pptx.oxml.ns import qn
        from lxml import etree
        etree.SubElement(bodyPr, qn('a:noAutofit'))
    except Exception:
        pass
    # Split on \n for reliable line breaks (separate paragraphs)
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        if lsp:
            p.line_spacing = Pt(lsp)
    return box


def card(slide, left, top, width, height, fill=CARD_BG):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def accent_bar(slide, left, top, width, color):
    b = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.06))
    b.fill.solid()
    b.fill.fore_color.rgb = color
    b.line.fill.background()
    return b


def circle(slide, left, top, size, color):
    c = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, Inches(size), Inches(size))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    return c


def dashboard_ph(slide, left, top, width, height, label):
    """Dashboard screenshot placeholder."""
    card(slide, left, top, width, height, RGBColor(0x1E, 0x2D, 0x3D))
    accent_bar(slide, left, top, width, SKY_BLUE)
    tb(slide, left + Inches(0.5),
       top + height / 2 - Inches(0.5),
       width - Inches(1), Inches(1.0),
       f"[ DASHBOARD SCREENSHOT ]\n{label}",
       16, bold=True, color=MUTED, align=PP_ALIGN.CENTER)


def stat_card(slide, left, top, label, value, color=SKY_BLUE,
              w=Inches(2.85), h=Inches(1.65)):
    card(slide, left, top, w, h)
    tb(slide, left + Inches(0.15), top + Inches(0.12),
       w - Inches(0.3), Inches(0.65), label, 15, bold=True, color=MUTED)
    tb(slide, left + Inches(0.15), top + Inches(0.8),
       w - Inches(0.3), Inches(0.65), str(value), 30, bold=True, color=color)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ═══════════════════════════════════════════════════════════════════════
#  BUILD
# ═══════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ── SLIDE 1: TITLE ───────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)
accent_bar(sl, Inches(0), Inches(0), SLIDE_W, SKY_BLUE)

tb(sl, Inches(1.5), Inches(1.6), Inches(10.3), Inches(1),
   "Burn Equity Index", 54, bold=True, color=SKY_BLUE, align=PP_ALIGN.CENTER)
tb(sl, Inches(1.5), Inches(2.8), Inches(10.3), Inches(0.7),
   "Measuring Structural Burn-Care Access Inequity Across America",
   26, color=TEXT_CLR, align=PP_ALIGN.CENTER)
tb(sl, Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.5),
   "Challenge Area 3: Equitable Access", 22, bold=True,
   color=AMBER, align=PP_ALIGN.CENTER)
tb(sl, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.5),
   "Team 15  |  Jacob Tang  \u00b7  Madeline Rose Johnson  \u00b7  "
   "Yisihaq Yemiru  \u00b7  Mashfika",
   18, color=MUTED, align=PP_ALIGN.CENTER)
tb(sl, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.4),
   "HeatMap Hackathon  |  American Burn Association  \u00b7  BData",
   16, color=MUTED, align=PP_ALIGN.CENTER)

notes(sl, """[MADELINE] [0:00 - 0:25] 25s

Good morning everyone. Picture this -- a child reaches for a pot on the stove. In that split second, everything changes. The question that matters most is: how quickly can that child reach specialized burn care? Today we'll show you that the answer depends almost entirely on where you live. I'm Madeline, and together with Jacob, we're presenting the Burn Equity Index.""")


# ── SLIDE 2: THE PROBLEM ────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)

tb(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
   "The Problem: Access Is Not Equal", 38, bold=True, color=TEXT_CLR)
tb(sl, Inches(0.7), Inches(0.9), Inches(12), Inches(0.4),
   "Burn care is one of the most regionalized specialties in medicine",
   20, color=MUTED)

stat_card(sl, Inches(0.7), Inches(1.6), "Burn Centers (US)", "136", SKY_BLUE)
stat_card(sl, Inches(3.8), Inches(1.6), "Total Burn Beds", "2,080", SKY_BLUE)
stat_card(sl, Inches(6.9), Inches(1.6), "US Population", "331M", AMBER)
stat_card(sl, Inches(10.0), Inches(1.6), "States with No\nBurn Center", "8", VERMILLION)

# Hennepin card
card(sl, Inches(0.7), Inches(3.4), Inches(5.8), Inches(3.5))
tb(sl, Inches(1.0), Inches(3.55), Inches(5.3), Inches(0.4),
   "Hennepin County, MN", 24, bold=True, color=TEAL)
tb(sl, Inches(1.0), Inches(4.0), Inches(5.3), Inches(0.35),
   "Pop: 1.27 million  |  BEI: 0", 20, color=TEXT_CLR)
tb(sl, Inches(1.0), Inches(4.5), Inches(5.3), Inches(1.8),
   "10 min to a verified burn center\n"
   "Direct access, full pediatric capability\n"
   "Multiple facilities within 30 min",
   20, color=MUTED, lsp=32)

# Kittson card
card(sl, Inches(6.8), Inches(3.4), Inches(5.8), Inches(3.5))
tb(sl, Inches(7.1), Inches(3.55), Inches(5.3), Inches(0.4),
   "Kittson County, MN", 24, bold=True, color=VERMILLION)
tb(sl, Inches(7.1), Inches(4.0), Inches(5.3), Inches(0.35),
   "Pop: 4,191  |  BEI: 99.6", 20, color=TEXT_CLR)
tb(sl, Inches(7.1), Inches(4.5), Inches(5.3), Inches(1.8),
   "8+ hours to nearest burn center\n"
   "No stabilization within 30 min\n"
   "Zero pediatric burn access",
   20, color=MUTED, lsp=32)

notes(sl, """[MADELINE] [0:25 - 1:10] 45s

There are only 136 verified burn centers and about 2,000 burn beds for 331 million Americans. Eight states have no burn center at all.

But the numbers only tell part of the story. In Hennepin County -- that's Minneapolis -- a burn patient reaches world-class care in 10 minutes. BEI score: zero. Perfect access. Now jump to Kittson County in the far northwest corner of the same state. Same Minnesota. But here, a burn patient faces an 8-hour journey, and the BEI is 99.6 -- nearly the worst score possible. That's the gap we set out to measure. Jacob, take us through how we did it.""")


# ── SLIDE 3: USE CASE ────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)

tb(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
   "Our Use Case: Equitable Access", 38, bold=True, color=TEXT_CLR)
tb(sl, Inches(0.7), Inches(1.0), Inches(11.5), Inches(1.1),
   "Given where people live and where burn-care resources sit, "
   "how equitably is the system structured?",
   24, color=TEXT_CLR, lsp=34)

pillars = [
    ("Quantify", "Build a composite equity index at\n"
     "census-tract resolution across\nthe entire United States.", SKY_BLUE),
    ("Visualize", "Create an interactive dashboard\n"
     "for stakeholders to explore access\nfrom national to neighborhood level.", AMBER),
    ("Act", "Identify high-burden communities\n"
     "and provide actionable, evidence-\n"
     "based recommendations.", TEAL),
]
for i, (title, desc, clr) in enumerate(pillars):
    x = Inches(0.7) + i * Inches(4.1)
    y = Inches(2.3)
    card(sl, x, y, Inches(3.8), Inches(4.3))
    circle(sl, x + Inches(1.45), y + Inches(0.3), 0.8, clr)
    tb(sl, x + Inches(1.45), y + Inches(0.35), Inches(0.8), Inches(0.7),
       str(i + 1), 28, bold=True, color=DARK_BG2, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.3), y + Inches(1.3), Inches(3.2), Inches(0.5),
       title, 26, bold=True, color=clr, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.3), y + Inches(1.9), Inches(3.2), Inches(2.0),
       desc, 18, color=MUTED, align=PP_ALIGN.CENTER, lsp=28)

notes(sl, """[JACOB] [1:10 - 1:45] 35s

Thanks, Madeline. I'm Jacob, and I led the technical build. Our primary use case is Equitable Access -- we're not predicting who gets burned or rating hospital quality. We're mapping structural inequity.

Three goals. First, quantify access with a rigorous index at census-tract resolution. Second, visualize it through an interactive dashboard anyone can use. Third, turn those findings into specific, actionable recommendations. Let me show you how the index works.""")


# ── SLIDE 4: THE BEI ─────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)

tb(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
   "The Burn Equity Index", 38, bold=True, color=TEXT_CLR)

tb(sl, Inches(0.7), Inches(1.2), Inches(12), Inches(0.7),
   "BEI  =  100  \u00d7  ( 0.25\u00b7S  +  0.30\u00b7T  +  0.20\u00b7P  +  0.25\u00b7C )",
   28, bold=True, color=AMBER, align=PP_ALIGN.CENTER, font="Consolas")

tb(sl, Inches(0.7), Inches(2.0), Inches(12), Inches(0.4),
   "Four pillars -- each captures a distinct dimension of structural access.",
   20, color=MUTED, align=PP_ALIGN.CENTER)

pillar_data = [
    ("S", "Supply\nScarcity", "25%",
     "How scarce is specialized\nburn-care supply nearby?", PILLAR_S),
    ("T", "Timely\nAccess", "30%",
     "How long to reach\ndefinitive burn care?", PILLAR_T),
    ("P", "Pediatric\nAccess", "20%",
     "Can children reach\npediatric burn care?", PILLAR_P),
    ("C", "Capacity\nGap", "25%",
     "Are there enough\nburn beds nearby?", PILLAR_C),
]
for i, (icon, name, weight, desc, clr) in enumerate(pillar_data):
    x = Inches(0.5) + i * Inches(3.15)
    y = Inches(2.7)
    card(sl, x, y, Inches(2.95), Inches(3.9))
    accent_bar(sl, x, y, Inches(2.95), clr)
    circle(sl, x + Inches(1.0), y + Inches(0.25), 0.85, clr)
    tb(sl, x + Inches(1.0), y + Inches(0.32), Inches(0.85), Inches(0.7),
       icon, 30, bold=True, color=DARK_BG2, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.2), y + Inches(1.3), Inches(2.55), Inches(0.85),
       name, 22, bold=True, color=TEXT_CLR, align=PP_ALIGN.CENTER, lsp=28)
    tb(sl, x + Inches(0.2), y + Inches(2.2), Inches(2.55), Inches(0.35),
       f"Weight: {weight}", 17, bold=True, color=clr, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.2), y + Inches(2.65), Inches(2.55), Inches(1.0),
       desc, 17, color=MUTED, align=PP_ALIGN.CENTER, lsp=26)

# Legend
tb(sl, Inches(1.5), Inches(6.8), Inches(3), Inches(0.3),
   "0 = Best Access", 17, bold=True, color=TEAL)
tb(sl, Inches(9.5), Inches(6.8), Inches(3), Inches(0.3),
   "100 = Worst Access", 17, bold=True, color=VERMILLION,
   align=PP_ALIGN.RIGHT)

notes(sl, """[JACOB] [1:45 - 2:25] 40s

The BEI scores every census tract in the country from zero to a hundred. Zero means excellent structural access; a hundred means severe barriers.

It blends four pillars. S captures how scarce specialized supply is near a community. T -- our heaviest weight at 30 percent -- measures how long it takes to reach definitive care, including transfer pathways. P isolates pediatric access, because children need different facilities. And C looks at whether there are enough burn beds for the surrounding population. Each one is a distinct lens on the same system. Let me walk through the key ones.""")


# ── SLIDE 5: PILLAR S ────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)
accent_bar(sl, Inches(0), Inches(0), SLIDE_W, PILLAR_S)

tb(sl, Inches(0.7), Inches(0.3), Inches(1), Inches(0.6),
   "S", 34, bold=True, color=PILLAR_S)
tb(sl, Inches(1.6), Inches(0.3), Inches(10), Inches(0.6),
   "Specialized Supply Scarcity  (25%)", 34, bold=True, color=TEXT_CLR)

tb(sl, Inches(0.7), Inches(1.1), Inches(11.5), Inches(0.5),
   "How scarce is specialized burn-care supply around each community?",
   22, color=MUTED)

# Facility weights
tb(sl, Inches(0.7), Inches(1.8), Inches(6), Inches(0.4),
   "Facility Capability Weights", 22, bold=True, color=PILLAR_S)
weights = [
    ("ABA-Verified Burn Center", "1.00", TEAL),
    ("State-Designated Burn Center", "0.85", SKY_BLUE),
    ("Burn-Capable (Non-Verified)", "0.50", AMBER),
    ("Trauma-Only Stabilization", "0.20", MUTED),
]
for j, (fac, wt, clr) in enumerate(weights):
    y = Inches(2.3) + j * Inches(0.5)
    tb(sl, Inches(1.0), y, Inches(5), Inches(0.4), fac, 18, color=TEXT_CLR)
    tb(sl, Inches(5.5), y, Inches(1.2), Inches(0.4),
       wt, 18, bold=True, color=clr, align=PP_ALIGN.RIGHT)

# How It Works
card(sl, Inches(7.3), Inches(1.8), Inches(5.5), Inches(4.6))
tb(sl, Inches(7.6), Inches(1.95), Inches(5), Inches(0.4),
   "How It Works (E2SFCA)", 20, bold=True, color=PILLAR_S)
steps_text = [
    "1. Compute each facility's supply-to-\n"
    "    demand ratio (capability weight\n"
    "    divided by nearby population).",
    "2. For each tract, sum ratios of\n"
    "    all reachable facilities.",
    "3. Normalize nationally, then flip:\n"
    "    high value = scarce supply.",
]
yp = Inches(2.5)
for st in steps_text:
    tb(sl, Inches(7.6), yp, Inches(4.9), Inches(1.1),
       st, 17, color=TEXT_CLR, lsp=24)
    yp += Inches(1.15)

# Distance decay
tb(sl, Inches(0.7), Inches(4.6), Inches(6), Inches(0.4),
   "Distance Decay Function", 22, bold=True, color=TEXT_CLR)
decay = [
    ("\u2264 30 min", "1.00  Full access"),
    ("31\u201360 min", "0.60  Reduced"),
    ("61\u201390 min", "0.30  Marginal"),
    ("> 90 min", "0.00  Out of reach"),
]
for j, (t, d) in enumerate(decay):
    y = Inches(5.1) + j * Inches(0.45)
    tb(sl, Inches(1.0), y, Inches(1.5), Inches(0.35),
       t, 17, bold=True, color=SKY_BLUE)
    tb(sl, Inches(2.7), y, Inches(3.5), Inches(0.35), d, 17, color=MUTED)

notes(sl, """[JACOB] [2:25 - 2:55] 30s

Supply uses a two-step floating catchment area method. It's a well-established spatial accessibility framework, and here's how it works in plain terms.

Each facility gets a capability weight -- a verified burn center is 1.0, a trauma-only site is 0.20. We then look at how many people live within driving range, using a step-decay function: full weight within 30 minutes, partial out to 90, and zero beyond. For each tract, we add up the ratios of all hospitals it can reach. Then we normalize and flip, so a high score means scarce supply.""")


# ── SLIDE 6: PILLAR T ────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)
accent_bar(sl, Inches(0), Inches(0), SLIDE_W, PILLAR_T)

tb(sl, Inches(0.7), Inches(0.3), Inches(1), Inches(0.6),
   "T", 34, bold=True, color=PILLAR_T)
tb(sl, Inches(1.6), Inches(0.3), Inches(10), Inches(0.6),
   "Timely Access Burden  (30% -- Highest Weight)", 34, bold=True,
   color=TEXT_CLR)

tb(sl, Inches(0.7), Inches(1.1), Inches(11.5), Inches(0.55),
   "Models the real-world tiered burn care system, not just distance to the nearest center.",
   20, color=MUTED, lsp=28)

# Three pathway cards
paths = [
    ("A", "Direct Path", "T(dir)",
     "Drive from the tract to the\nnearest definitive burn center.", SKY_BLUE),
    ("B", "Transfer Path", "T(trans)",
     "Drive to stabilization hospital,\nadd 45-min transfer penalty,\n"
     "then transfer to burn center.", AMBER),
    ("C", "System Time", "T(sys)",
     "Whichever is faster -- direct\n"
     "or stabilize-then-transfer.", TEAL),
]
for i, (letter, title, formula, desc, clr) in enumerate(paths):
    x = Inches(0.5) + i * Inches(4.15)
    y = Inches(1.8)
    card(sl, x, y, Inches(3.9), Inches(2.7))
    circle(sl, x + Inches(0.2), y + Inches(0.2), 0.6, clr)
    tb(sl, x + Inches(0.2), y + Inches(0.22), Inches(0.6), Inches(0.55),
       letter, 24, bold=True, color=DARK_BG2, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(1.0), y + Inches(0.2), Inches(2.6), Inches(0.4),
       title, 20, bold=True, color=TEXT_CLR)
    tb(sl, x + Inches(1.0), y + Inches(0.6), Inches(2.6), Inches(0.3),
       formula, 15, bold=True, color=clr, font="Consolas")
    tb(sl, x + Inches(0.25), y + Inches(1.1), Inches(3.4), Inches(1.3),
       desc, 17, color=MUTED, lsp=26)

# Tier penalty
card(sl, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.1))
tb(sl, Inches(0.8), Inches(4.95), Inches(11.7), Inches(0.4),
   "Tier Penalty:  \u0394 = max(0, T_stab \u2212 30 min)", 20, bold=True,
   color=VERMILLION, font="Consolas")
tb(sl, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.2),
   "We also penalize tracts where even first-line stabilization is more "
   "than 30 minutes away. Final score: 75% system time + 25% tier penalty. "
   "This captures both the journey and the safety net.",
   19, color=TEXT_CLR, lsp=30)

notes(sl, """[JACOB] [2:55 - 3:45] 50s

This is our heaviest-weighted pillar, because in burn care, minutes matter. And what sets our approach apart is that we model how the system actually works.

We compute three pathways. Path A is direct -- just the drive time to the nearest burn center. Path B is the transfer route -- go to the nearest stabilization hospital, add a 45-minute structural transfer penalty, then continue to a definitive center. Path C picks whichever is faster.

On top of that, we add a tier penalty. If even your nearest emergency room is more than 30 minutes away, that compounds the problem. So the final score blends 75 percent system travel time with 25 percent tier penalty -- capturing both the journey and the safety net.""")


# ── SLIDE 7: PILLARS P & C ──────────────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)

# Left half: P
accent_bar(sl, Inches(0), Inches(0), Inches(6.5), PILLAR_P)
tb(sl, Inches(0.7), Inches(0.3), Inches(0.8), Inches(0.6),
   "P", 34, bold=True, color=PILLAR_P)
tb(sl, Inches(1.5), Inches(0.3), Inches(5), Inches(0.6),
   "Pediatric Access Gap  (20%)", 30, bold=True, color=TEXT_CLR)

tb(sl, Inches(0.7), Inches(1.1), Inches(5.6), Inches(1.1),
   "Same method as Supply, but using\n"
   "pediatric-capable facilities and\nchild population.",
   18, color=MUTED, lsp=26)

peds = [
    ("Pediatric ABA-Verified", "1.00"),
    ("Pediatric State-Designated", "0.85"),
    ("Peds Trauma + Burn Capable", "0.60"),
    ("Pediatric Stabilization Only", "0.25"),
]
for j, (fac, wt) in enumerate(peds):
    y = Inches(2.3) + j * Inches(0.45)
    tb(sl, Inches(0.9), y, Inches(4.3), Inches(0.35), fac, 17, color=TEXT_CLR)
    tb(sl, Inches(5.0), y, Inches(1.2), Inches(0.35),
       wt, 17, bold=True, color=PILLAR_P, align=PP_ALIGN.RIGHT)

card(sl, Inches(0.7), Inches(4.5), Inches(5.6), Inches(2.2))
tb(sl, Inches(0.9), Inches(4.65), Inches(5.2), Inches(0.35),
   "Why a separate pillar?", 19, bold=True, color=PILLAR_P)
tb(sl, Inches(0.9), Inches(5.1), Inches(5.2), Inches(1.4),
   "Pediatric burns need different expertise\n"
   "and equipment. A community can have\n"
   "adult access but zero pediatric capability.",
   17, color=TEXT_CLR, lsp=26)

# Right half: C
accent_bar(sl, Inches(6.8), Inches(0), Inches(6.5), PILLAR_C)
tb(sl, Inches(7.0), Inches(0.3), Inches(0.8), Inches(0.6),
   "C", 34, bold=True, color=PILLAR_C)
tb(sl, Inches(7.8), Inches(0.3), Inches(5.2), Inches(0.6),
   "Structural Capacity Gap  (25%)", 30, bold=True, color=TEXT_CLR)

tb(sl, Inches(7.0), Inches(1.1), Inches(5.6), Inches(1.1),
   "Measures burn-bed adequacy relative\n"
   "to surrounding population, weighted\nby distance.",
   18, color=MUTED, lsp=26)

card(sl, Inches(7.0), Inches(2.2), Inches(5.8), Inches(2.0))
tb(sl, Inches(7.3), Inches(2.35), Inches(5.2), Inches(0.35),
   "Effective Beds = BURN_BEDS \u00d7 u", 18, bold=True,
   color=PILLAR_C, font="Consolas")
tb(sl, Inches(7.3), Inches(2.8), Inches(5.2), Inches(1.1),
   "u = capacity utilization factor\n"
   "Baseline: u = 1.0 (structural)\n"
   "Sensitivity: u = 0.75 (conservative)",
   17, color=TEXT_CLR, lsp=26)

card(sl, Inches(7.0), Inches(4.5), Inches(5.8), Inches(2.2))
tb(sl, Inches(7.3), Inches(4.65), Inches(5.2), Inches(0.35),
   "The national picture", 19, bold=True, color=PILLAR_C)
tb(sl, Inches(7.3), Inches(5.1), Inches(5.2), Inches(1.4),
   "2,080 total US burn beds\n"
   "= 0.63 beds per 100,000 people\n"
   "Concentrated in urban corridors",
   17, color=TEXT_CLR, lsp=26)

notes(sl, """[JACOB] [3:45 - 4:20] 35s

The Pediatric pillar uses the same framework but calibrated for children -- weighting facilities by pediatric burn capability and using child population as the denominator. We gave it its own pillar because a community might have decent adult access but nothing for kids.

Capacity measures bed adequacy. Across the whole US, there are just 2,080 burn beds -- that's 0.63 per hundred thousand people. And those beds cluster in urban centers, so rural communities are doubly disadvantaged: they're both far away and under-bedded.""")


# ── SLIDE 8: DATA SOURCES ────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)

tb(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
   "Data Sources & Integration", 38, bold=True, color=TEXT_CLR)
tb(sl, Inches(0.7), Inches(0.9), Inches(12), Inches(0.4),
   "Every source is public or challenge-provided -- fully reproducible",
   20, color=MUTED)

sources = [
    ("NIRD", "635 facilities with burn\ndesignation and bed counts", SKY_BLUE),
    ("US Census (ACS)", "84,000+ tract boundaries\nand population data", AMBER),
    ("OpenStreetMap\n+ Valhalla", "Real road-network travel\ntimes, not estimates", TEAL),
    ("RUCA Codes", "Rural-urban classification\nfor every census tract", PINK),
    ("CDC SVI", "Social vulnerability overlay\n(kept outside core BEI)", MUTED),
    ("FAA Airport Data", "Heliport locations for air-\ntransport sensitivity", VERMILLION),
]
for i, (title, desc, clr) in enumerate(sources):
    c = i % 3
    r = i // 3
    x = Inches(0.5) + c * Inches(4.15)
    y = Inches(1.6) + r * Inches(2.6)
    card(sl, x, y, Inches(3.9), Inches(2.3))
    tb(sl, x + Inches(0.25), y + Inches(0.15), Inches(3.4), Inches(0.7),
       title, 19, bold=True, color=clr, lsp=24)
    tb(sl, x + Inches(0.25), y + Inches(0.9), Inches(3.4), Inches(1.1),
       desc, 17, color=TEXT_CLR, lsp=24)

notes(sl, """[JACOB] [4:20 - 4:45] 25s

Every data source is either provided by the challenge or freely available. NIRD gives us 635 hospitals. Census provides tract-level population. We compute actual drive times using OpenStreetMap and a local Valhalla routing engine -- these are real road-network times, not straight-line guesses. We also layer in rural-urban classifications, social vulnerability data, and FAA airport records for our air-transport scenario. The entire pipeline is reproducible.""")


# ── SLIDE 9: TRANSITION ──────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)

tb(sl, Inches(1.5), Inches(2.0), Inches(10.3), Inches(0.9),
   "What Does the Data Reveal?", 46, bold=True, color=TEXT_CLR,
   align=PP_ALIGN.CENTER)
tb(sl, Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.5),
   "From method to map -- let's explore what the data shows",
   24, color=AMBER, align=PP_ALIGN.CENTER)

stat_card(sl, Inches(0.8), Inches(4.4), "MN Tracts", "1,505", SKY_BLUE)
stat_card(sl, Inches(3.95), Inches(4.4), "MN Median BEI", "23.0", AMBER)
stat_card(sl, Inches(7.1), Inches(4.4), "High-Burden\n(BEI \u2265 80)", "155 tracts", VERMILLION)
stat_card(sl, Inches(10.25), Inches(4.4), "US Counties", "3,144", SKY_BLUE)

notes(sl, """[MADELINE] [4:45 - 5:05] 20s

That's the method. Now let's see what it finds. We scored 1,505 tracts across Minnesota and 3,144 counties nationwide. Minnesota's median BEI is 23 -- that sounds good. But 155 tracts score above 80. Those are communities where reaching burn care takes hours, not minutes. Let me show you on the map.""")


# ── SLIDE 10: MN MAP (DASHBOARD) ─────────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)

tb(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
   "Minnesota: Burn Equity Index Map", 34, bold=True, color=TEXT_CLR)
tb(sl, Inches(0.7), Inches(0.85), Inches(12), Inches(0.35),
   "Tract-level BEI scores  |  Interactive Dashboard", 18, color=MUTED)

dashboard_ph(sl, Inches(0.7), Inches(1.4), Inches(8.3), Inches(5.5),
             "MN Map Page -- BEI choropleth with facility markers")

# Side callouts
card(sl, Inches(9.3), Inches(1.4), Inches(3.5), Inches(2.5))
tb(sl, Inches(9.5), Inches(1.55), Inches(3.1), Inches(0.35),
   "Twin Cities Metro", 19, bold=True, color=TEAL)
tb(sl, Inches(9.5), Inches(1.95), Inches(3.1), Inches(1.6),
   "BEI: 0 -- 2\n"
   "Multiple burn centers\n"
   "within 15 minutes\n"
   "50%+ of MN population",
   17, color=TEXT_CLR, lsp=26)

card(sl, Inches(9.3), Inches(4.2), Inches(3.5), Inches(2.7))
tb(sl, Inches(9.5), Inches(4.35), Inches(3.1), Inches(0.35),
   "Northwest MN", 19, bold=True, color=VERMILLION)
tb(sl, Inches(9.5), Inches(4.75), Inches(3.1), Inches(1.8),
   "BEI: 90 -- 100\n"
   "No burn center for\nhundreds of miles\n"
   "8+ hour travel time",
   17, color=TEXT_CLR, lsp=26)

notes(sl, """[MADELINE] [5:05 - 5:40] 35s

[USE LIVE DASHBOARD OR SCREENSHOT]

Here's Minnesota, colored by BEI. The Twin Cities metro shows green -- near-zero scores. Half the state's population lives within minutes of a burn center. Now look at the northwest corner. Deep red. Scores above 90. Kittson, Roseau, Marshall counties. No burn center for hundreds of miles. The transition isn't gradual. It drops off like a cliff.""")


# ── SLIDE 11: COUNTY STORIES ─────────────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)

tb(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
   "Minnesota County Stories", 34, bold=True, color=TEXT_CLR)
tb(sl, Inches(0.7), Inches(0.85), Inches(12), Inches(0.35),
   "Same state, vastly different realities", 20, color=MUTED)

counties = [
    ("Hennepin County", "Minneapolis", "BEI: 0.0",
     "10.6 min to burn center\nDirect access pathway\n"
     "All 4 pillars score 0.00\nPop: 1.27 million",
     TEAL, "Best Access"),
    ("Blue Earth County", "Mankato", "BEI: 58.5",
     "88 min to burn center\nS: 1.00  T: 0.14  C: 1.00\n"
     "Supply & capacity gaps\nPop: 69,022",
     AMBER, "Moderate Burden"),
    ("Beltrami County", "Bemidji", "BEI: 90.4",
     "8+ hrs to burn center\nTransfer pathway only\n"
     "S, T, and C all maxed\nPop: 46,274",
     VERMILLION, "Severe Burden"),
]
for i, (county, city, bei, details, clr, label) in enumerate(counties):
    x = Inches(0.5) + i * Inches(4.15)
    y = Inches(1.5)
    card(sl, x, y, Inches(3.9), Inches(5.4))
    accent_bar(sl, x, y, Inches(3.9), clr)
    tb(sl, x + Inches(0.25), y + Inches(0.25), Inches(3.4), Inches(0.3),
       label, 15, bold=True, color=clr)
    tb(sl, x + Inches(0.25), y + Inches(0.6), Inches(3.4), Inches(0.45),
       county, 22, bold=True, color=TEXT_CLR)
    tb(sl, x + Inches(0.25), y + Inches(1.05), Inches(3.4), Inches(0.3),
       city, 17, color=MUTED)
    tb(sl, x + Inches(0.25), y + Inches(1.5), Inches(3.4), Inches(0.5),
       bei, 26, bold=True, color=clr)
    tb(sl, x + Inches(0.25), y + Inches(2.3), Inches(3.4), Inches(2.3),
       details, 17, color=TEXT_CLR, lsp=30)

notes(sl, """[MADELINE] [5:40 - 6:25] 45s

Three counties, one state.

Hennepin -- Minneapolis. BEI: zero. You're 10 minutes from definitive care. Every pillar scores zero. That's what good access looks like.

Blue Earth -- Mankato. A university town, about 69,000 people. BEI: 58.5. You're an hour and a half from a burn center by road. The supply and capacity pillars are both maxed out -- there's simply nothing nearby. But the timely-access score is low because the drive is technically still feasible. A community like this can look fine on a map but is structurally underserved.

Beltrami -- Bemidji. BEI: 90.4. Eight hours to a burn center. Transfer pathway only. Supply, time, and capacity all maxed out. That is structural isolation.""")


# ── SLIDE 12: RURAL-URBAN (DASHBOARD) ────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)

tb(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
   "The Rural\u2013Urban Divide", 34, bold=True, color=TEXT_CLR)

dashboard_ph(sl, Inches(0.5), Inches(1.1), Inches(7.2), Inches(5.8),
             "MN rural vs urban box plot\n(mn_03_rural_urban_gap.png)")

# Right stats
card(sl, Inches(8.0), Inches(1.1), Inches(4.8), Inches(2.2))
tb(sl, Inches(8.3), Inches(1.25), Inches(4.2), Inches(0.4),
   "Median Travel Time", 20, bold=True, color=TEXT_CLR)
tb(sl, Inches(8.3), Inches(1.75), Inches(2.0), Inches(0.4),
   "Urban:", 18, color=MUTED)
tb(sl, Inches(10.2), Inches(1.75), Inches(2.3), Inches(0.4),
   "19 min", 22, bold=True, color=TEAL)
tb(sl, Inches(8.3), Inches(2.25), Inches(2.0), Inches(0.4),
   "Rural:", 18, color=MUTED)
tb(sl, Inches(10.2), Inches(2.25), Inches(2.3), Inches(0.4),
   "124 min", 22, bold=True, color=VERMILLION)

card(sl, Inches(8.0), Inches(3.6), Inches(4.8), Inches(1.5))
tb(sl, Inches(8.3), Inches(3.75), Inches(4.2), Inches(0.5),
   "6.5x Longer", 26, bold=True, color=VERMILLION)
tb(sl, Inches(8.3), Inches(4.3), Inches(4.2), Inches(0.75),
   "Rural residents travel 6.5 times\nfarther to reach burn care.",
   18, color=TEXT_CLR, lsp=26)

card(sl, Inches(8.0), Inches(5.3), Inches(4.8), Inches(1.5))
tb(sl, Inches(8.3), Inches(5.45), Inches(4.2), Inches(1.25),
   "30.5% of MN tracts are rural\n"
   "Only 50.1% of Minnesotans\n"
   "can reach care in 30 min",
   18, color=MUTED, lsp=28)

notes(sl, """[MADELINE] [6:25 - 6:55] 30s

This chart puts the disparity in sharp relief. Urban Minnesotans: 19-minute median. Rural Minnesotans: 124 minutes -- over two hours, and that's the median, so half face even longer trips. A 6.5x gap within one state. And nearly a third of the state's tracts are classified as rural. For burns, where that first hour can determine whether a patient needs grafts or makes a full recovery, two hours is not just inconvenient. It changes outcomes.""")


# ── SLIDE 13: NATIONAL MAP (DASHBOARD) ───────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)

tb(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
   "National: County-Level BEI Map", 34, bold=True, color=TEXT_CLR)
tb(sl, Inches(0.7), Inches(0.85), Inches(12), Inches(0.35),
   "3,144 counties  |  331 million Americans  |  Interactive Dashboard",
   18, color=MUTED)

dashboard_ph(sl, Inches(0.5), Inches(1.4), Inches(8.5), Inches(5.5),
             "USA County Map Page -- Full national BEI choropleth")

card(sl, Inches(9.3), Inches(1.4), Inches(3.5), Inches(1.7))
tb(sl, Inches(9.5), Inches(1.5), Inches(3.1), Inches(0.35),
   "Worst States (Avg BEI)", 17, bold=True, color=VERMILLION)
tb(sl, Inches(9.5), Inches(1.9), Inches(3.1), Inches(1.0),
   "AK: 99.7   ND: 99.0\n"
   "MT: 96.5   SD: 96.4\n"
   "WY: 83.4   NV: 83.3",
   16, color=TEXT_CLR, lsp=24)

card(sl, Inches(9.3), Inches(3.3), Inches(3.5), Inches(1.4))
tb(sl, Inches(9.5), Inches(3.4), Inches(3.1), Inches(0.35),
   "Best States (Avg BEI)", 17, bold=True, color=TEAL)
tb(sl, Inches(9.5), Inches(3.8), Inches(3.1), Inches(0.7),
   "NJ: 10.0   RI: 29.5\n"
   "MD: 29.8   MA: 31.0",
   16, color=TEXT_CLR, lsp=24)

card(sl, Inches(9.3), Inches(4.9), Inches(3.5), Inches(2.0))
tb(sl, Inches(9.5), Inches(5.0), Inches(3.1), Inches(0.6),
   "8 States -- Zero Burn\nCenters", 17, bold=True, color=AMBER, lsp=22)
tb(sl, Inches(9.5), Inches(5.65), Inches(3.1), Inches(0.85),
   "AK, DE, MS, MT,\nND, NH, SD, WY",
   16, color=TEXT_CLR, lsp=24)

notes(sl, """[MADELINE] [6:55 - 7:30] 35s

[SWITCH TO NATIONAL DASHBOARD VIEW IF POSSIBLE]

Now the full country. The coasts light up green -- New Jersey, Maryland, Massachusetts, all below 31. The interior is a different story. Great Plains, Mountain West, rural Deep South -- all high-burden. Alaska and North Dakota average above 99. Eight states have no burn center at all. If you're burned in Montana, you have to leave the state. These aren't empty places -- Flathead County, Montana: 106,000 residents, BEI of 100.""")


# ── SLIDE 14: STATE RANKINGS (DASHBOARD) ─────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)

tb(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
   "State Rankings & Coverage Gaps", 34, bold=True, color=TEXT_CLR)

dashboard_ph(sl, Inches(0.4), Inches(1.2), Inches(6.1), Inches(5.7),
             "State Rankings bar chart\n(usa_02_state_rankings.png)")
dashboard_ph(sl, Inches(6.7), Inches(1.2), Inches(6.1), Inches(5.7),
             "National Coverage Gap chart\n(usa_03_coverage_gap.png)")

notes(sl, """[JACOB] [7:30 - 8:00] 30s

The state rankings confirm the geographic pattern. States with vast rural areas and few burn centers sit at the bottom. But population-weighted BEI also reveals inequity in states that seem well-resourced on paper. The coverage gap chart on the right shows millions of Americans -- disproportionately in rural and tribal communities -- living more than 60 minutes from verified care. That matters clinically. The ABA's own referral guidelines stress timely transfer for major burns.""")


# ── SLIDE 15: INNOVATION ─────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)

tb(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
   "What Makes BEI Different", 38, bold=True, color=TEXT_CLR)

inno = [
    ("Tiered System Modeling",
     "Models the real-world stabilize-then-\n"
     "transfer pathway with explicit transfer\n"
     "penalties and tier gap detection.",
     SKY_BLUE),
    ("Multi-Scenario Transport",
     "Ground-only baseline plus ground-\n"
     "plus-air sensitivity using FAA data.\n"
     "Published as scenario analysis.",
     AMBER),
    ("Four-Pillar Composite",
     "Supply, timeliness, pediatrics, and\n"
     "capacity each address a distinct\n"
     "clinical dimension. Fully transparent.",
     TEAL),
    ("Interactive Dashboard",
     "Any stakeholder can explore access\n"
     "patterns from national overview down\n"
     "to individual tracts. No coding needed.",
     VERMILLION),
]
for i, (title, desc, clr) in enumerate(inno):
    c = i % 2
    r = i // 2
    x = Inches(0.5) + c * Inches(6.3)
    y = Inches(1.2) + r * Inches(3.0)
    card(sl, x, y, Inches(6.0), Inches(2.7))
    circle(sl, x + Inches(0.2), y + Inches(0.2), 0.55, clr)
    tb(sl, x + Inches(0.2), y + Inches(0.23), Inches(0.55), Inches(0.5),
       str(i + 1), 20, bold=True, color=DARK_BG2, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(1.0), y + Inches(0.2), Inches(4.7), Inches(0.4),
       title, 21, bold=True, color=TEXT_CLR)
    tb(sl, x + Inches(0.25), y + Inches(0.8), Inches(5.5), Inches(1.5),
       desc, 17, color=MUTED, lsp=26)

notes(sl, """[JACOB] [8:00 - 8:35] 35s

Four things set the BEI apart. First, we model the real tiered care system -- not just nearest-facility distance, but the full stabilize-and-transfer pathway with explicit penalties.

Second, we run multi-scenario transport. Ground-only as the reproducible baseline, plus a ground-plus-air sensitivity scenario built on real FAA infrastructure data.

Third, four pillars that each address a different clinical dimension. And fourth, an interactive dashboard that makes all of this accessible to anyone -- clinicians, planners, or the public. No technical expertise needed.""")


# ── SLIDE 16: RECOMMENDATIONS ────────────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)

tb(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
   "Impact & Recommendations", 38, bold=True, color=TEXT_CLR)
tb(sl, Inches(0.7), Inches(0.9), Inches(12), Inches(0.35),
   "From data to decisions", 20, color=MUTED)

recs = [
    ("Telemedicine Triage",
     "Deploy tele-burn consultation in\n"
     "high-BEI regions. Research shows\n"
     "94% accuracy in remote surgical\n"
     "determination.",
     SKY_BLUE, "Referral"),
    ("Strategic Facility\nPlacement",
     "Use hotspot clusters to identify\n"
     "where new burn-capable facilities\n"
     "would have the greatest impact.\n"
     "NW Minnesota is a clear candidate.",
     AMBER, "Equitable Access"),
    ("Air Transport\nInvestment",
     "Target air-ambulance resources\n"
     "where ground BEI exceeds 80 but\n"
     "air scenario shows 5+ point gain.\n"
     "Existing heliports are underused.",
     TEAL, "Referral"),
    ("Policy Planning Tool",
     "Integrate BEI into state health\n"
     "department planning and ABA\n"
     "verification processes. Open-source\n"
     "and fully reproducible.",
     VERMILLION, "Equitable Access"),
]
for i, (title, desc, clr, uc) in enumerate(recs):
    c = i % 2
    r = i // 2
    x = Inches(0.5) + c * Inches(6.3)
    y = Inches(1.5) + r * Inches(2.8)
    card(sl, x, y, Inches(6.0), Inches(2.5))
    tb(sl, x + Inches(0.25), y + Inches(0.15), Inches(3.8), Inches(0.7),
       title, 20, bold=True, color=TEXT_CLR, lsp=24)
    tb(sl, x + Inches(4.2), y + Inches(0.15), Inches(1.5), Inches(0.35),
       uc, 13, bold=True, color=clr, align=PP_ALIGN.RIGHT)
    tb(sl, x + Inches(0.25), y + Inches(0.85), Inches(5.5), Inches(1.4),
       desc, 16, color=MUTED, lsp=24)

notes(sl, """[MADELINE] [8:35 - 9:20] 45s

So what do we do with this data? Four things.

Telemedicine. In places like Beltrami County, a tele-burn consult can guide treatment while the patient is in transit. Research shows 94 percent accuracy in remote surgical assessment.

Facility placement. Our hotspot analysis shows exactly where a new burn-capable site would close the biggest gaps. Northwest Minnesota is the clear candidate.

Air transport. Our analysis finds corridors where adding air access drops the BEI significantly. Existing heliports in those areas are underused.

And the BEI itself becomes the planning tool. Open-source, reproducible, and ready for state health departments and the ABA to adopt.""")


# ── SLIDE 17: LIMITATIONS ────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)

tb(sl, Inches(0.7), Inches(0.3), Inches(12), Inches(0.6),
   "Limitations & Future Directions", 34, bold=True, color=TEXT_CLR)

lims = [
    ("National routing at county level",
     "Masks variation within large rural counties. MN uses ideal tract-level routing. "
     "Full US tract-level needs a high-memory server."),
    ("Florida data gap",
     "67 FL counties absent due to a routing batch error. "
     "Pipeline issue, not methodological -- the formula applies identically."),
    ("Air scenario is structural",
     "An accessibility estimate, not an operational prediction. "
     "Real-time weather and dispatch are out of scope."),
    ("Static capacity model",
     "Uses structural bed counts, not real-time occupancy. "
     "Live bed-census integration is a clear next step."),
]
for i, (title, desc) in enumerate(lims):
    y = Inches(1.2) + i * Inches(1.4)
    card(sl, Inches(0.5), y, Inches(12.3), Inches(1.2))
    tb(sl, Inches(0.8), y + Inches(0.1), Inches(5), Inches(0.35),
       title, 19, bold=True, color=AMBER)
    tb(sl, Inches(0.8), y + Inches(0.5), Inches(11.5), Inches(0.55),
       desc, 17, color=TEXT_CLR)

tb(sl, Inches(0.7), Inches(6.9), Inches(12), Inches(0.3),
   "Next:  Tract-level national routing  |  Real-time bed feeds  |  "
   "Temporal BEI  |  State health planning integration",
   16, color=MUTED, align=PP_ALIGN.CENTER)

notes(sl, """[JACOB] [9:20 - 9:50] 30s

We want to be upfront about limitations. National analysis is at county level, which can mask variation in large rural counties -- Minnesota uses full tract-level routing. Florida's counties are absent due to a routing error in our pipeline, not a methodology gap. Our air scenario is a structural estimate, not an operational forecast. And bed capacity is structural, not real-time.

Going forward: tract-level national routing, live bed census feeds, temporal BEI that captures seasonal patterns, and direct integration with state planning systems.""")


# ── SLIDE 18: CLOSING ────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)

tb(sl, Inches(1.5), Inches(1.2), Inches(10.3), Inches(0.8),
   "Every Minute Matters", 46, bold=True, color=TEXT_CLR,
   align=PP_ALIGN.CENTER)

tb(sl, Inches(1.5), Inches(2.3), Inches(10.3), Inches(2.0),
   "The Burn Equity Index gives us a shared language\n"
   "to see, measure, and act on the structural gaps\n"
   "that separate 10 minutes from 10 hours.",
   26, color=MUTED, align=PP_ALIGN.CENTER, lsp=40)

stat_card(sl, Inches(1.5), Inches(4.7), "Hennepin County", "BEI: 0", TEAL)
stat_card(sl, Inches(5.2), Inches(4.7), "Same State", "10m vs 8h", AMBER)
stat_card(sl, Inches(8.9), Inches(4.7), "Kittson County", "BEI: 99.6", VERMILLION)

tb(sl, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.4),
   "The gap is structural. What's structural can be changed. "
   "Change starts with data.",
   21, bold=True, color=SKY_BLUE, align=PP_ALIGN.CENTER)

notes(sl, """[MADELINE] [9:50 - 10:45] 55s

Remember that child at the stove. In Hennepin County, they're in a burn unit in 10 minutes. In Kittson County -- same state -- eight hours.

We've scored every county in America and every census tract in Minnesota. The Burn Equity Index gives clinicians, planners, and policymakers one number that captures what used to take months to assess.

The gap is structural. But structural means it can be changed. A new facility, a telemedicine link, an air corridor -- each one shifts the score. And every score represents real families. That's what the BEI does.""")


# ── SLIDE 19: THANK YOU ──────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)
accent_bar(sl, Inches(0), Inches(0), SLIDE_W, SKY_BLUE)

tb(sl, Inches(1.5), Inches(0.8), Inches(10.3), Inches(0.7),
   "Thank You", 46, bold=True, color=TEXT_CLR, align=PP_ALIGN.CENTER)
tb(sl, Inches(1.5), Inches(1.6), Inches(10.3), Inches(0.4),
   "Team 15  \u2014  Burn Equity Index  \u2014  Challenge Area 3",
   20, color=MUTED, align=PP_ALIGN.CENTER)

team = [
    ("Jacob Tang", "Technical Lead\nPipeline & Methodology", SKY_BLUE),
    ("Madeline Rose Johnson", "Analytics & Presentation\nMS Data Science", AMBER),
    ("Yisihaq Yemiru", "Research &\nData Integration", TEAL),
    ("Mashfika", "Analysis &\nVisualization", PINK),
]
for i, (name, role, clr) in enumerate(team):
    x = Inches(0.8) + i * Inches(3.1)
    y = Inches(2.4)
    card(sl, x, y, Inches(2.8), Inches(2.3))
    circle(sl, x + Inches(0.95), y + Inches(0.2), 0.8, clr)
    initials = "".join(w[0] for w in name.split()[:2])
    tb(sl, x + Inches(0.95), y + Inches(0.27), Inches(0.8), Inches(0.65),
       initials, 22, bold=True, color=DARK_BG2, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), y + Inches(1.15), Inches(2.6), Inches(0.35),
       name, 17, bold=True, color=TEXT_CLR, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), y + Inches(1.5), Inches(2.6), Inches(0.55),
       role, 13, color=MUTED, align=PP_ALIGN.CENTER, lsp=18)

tb(sl, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.35),
   "American Burn Association  \u00b7  BData  \u00b7  HealthcareMN  \u00b7  "
   "MinneAnalytics  \u00b7  University of Minnesota",
   15, color=MUTED, align=PP_ALIGN.CENTER)

tb(sl, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.35),
   "Data: NIRD  \u00b7  US Census  \u00b7  OpenStreetMap  \u00b7  "
   "Valhalla  \u00b7  RUCA  \u00b7  CDC SVI  \u00b7  FAA",
   15, color=MUTED, align=PP_ALIGN.CENTER)

card(sl, Inches(3.5), Inches(6.1), Inches(6.3), Inches(0.6))
tb(sl, Inches(3.5), Inches(6.15), Inches(6.3), Inches(0.5),
   "Interactive Dashboard Available", 19, bold=True,
   color=SKY_BLUE, align=PP_ALIGN.CENTER)

notes(sl, """[BOTH] [10:45 - 10:59] 14s

MADELINE: Thank you for your time. We'd love to take questions, and the interactive dashboard is available for you to explore.

JACOB: The full pipeline, methodology, and dashboard are open and reproducible. Thank you.

[END -- TOTAL: 10:59]

=== SCRIPT SUMMARY ===
~1,600 words | 10:59

JACOB: Slides 3-8, 14-15, 17 (~5:10)
MADELINE: Slides 1-2, 9-13, 16, 18-19 (~5:49)

=== RUBRIC COVERAGE ===
I. Clinical/Business Use Case (15 pts)
   1. Use Case ID: Equitable Access (Slide 3)
   2. Insights: County stories, rural-urban gap (Slides 10-14)
   3. Impact: Recommendations with stakeholder relevance (Slide 16)

II. Analytic/Methodologic Quality (15 pts)
   4. Methods: E2SFCA, tiered routing, normalization (Slides 4-7)
   5. Innovation: Tiered system, multi-scenario, dashboard (Slide 15)
   6. Data Integration: 6 public/challenge sources (Slide 8)

III. Presentation & Communication (15 pts)
   7. Storytelling: Human hook, county stories, bookend close
   8. Visual Quality: Dashboard, CVD-safe colors, dual presenters
   9. Actionability: 4 concrete, feasible recommendations (Slide 16)
""")


# ═══════════════════════════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════════════════════════
prs.save(str(OUT))
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
